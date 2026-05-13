from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PATH = "Presentation_Template.pptx"

BG = RGBColor(0x11, 0x15, 0x1F)
PANEL = RGBColor(0x1A, 0x22, 0x30)
PANEL_ALT = RGBColor(0x20, 0x2B, 0x3D)
ACCENT = RGBColor(0xC8, 0xA9, 0x6B)
ACCENT_ALT = RGBColor(0x6F, 0x8F, 0xB8)
TEXT = RGBColor(0xF3, 0xF2, 0xEF)
MUTED = RGBColor(0xBF, 0xC6, 0xD1)
LINE = RGBColor(0x2C, 0x37, 0x49)

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_rect(slide, left, top, width, height, color, line_color=None, line_width=1, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = line_color or color
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color=TEXT,
    font_name=BODY_FONT,
    bold=False,
    align=PP_ALIGN.LEFT,
    vertical=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, title, bullets):
    add_textbox(slide, left, top, width, Inches(0.5), title, 20, color=TEXT, font_name=BODY_FONT, bold=True)
    box = slide.shapes.add_textbox(left, top + Inches(0.55), width, height - Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = MUTED
        p.space_after = Pt(10)
    return box


def add_header(slide, title, subtitle=None):
    add_line(slide, Inches(0.7), Inches(0.55), Inches(2.3), Inches(0.04), ACCENT)
    add_textbox(
        slide,
        Inches(0.7),
        Inches(0.72),
        Inches(8.6),
        Inches(0.6),
        title,
        28,
        color=TEXT,
        font_name=TITLE_FONT,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.7),
            Inches(1.2),
            Inches(8.8),
            Inches(0.45),
            subtitle,
            12,
            color=MUTED,
            font_name=BODY_FONT,
        )


def add_footer(slide, center_text="Course / Class", right_text="Presenter Name | Date"):
    add_line(slide, Inches(0.7), Inches(6.82), Inches(11.93), Inches(0.015), LINE)
    add_textbox(
        slide,
        Inches(0.72),
        Inches(6.88),
        Inches(3.0),
        Inches(0.25),
        center_text,
        10,
        color=MUTED,
        font_name=BODY_FONT,
    )
    add_textbox(
        slide,
        Inches(9.35),
        Inches(6.88),
        Inches(3.2),
        Inches(0.25),
        right_text,
        10,
        color=MUTED,
        font_name=BODY_FONT,
        align=PP_ALIGN.RIGHT,
    )


def add_corner_accent(slide):
    add_rect(slide, Inches(10.9), Inches(0.45), Inches(1.2), Inches(1.2), PANEL_ALT, line_color=ACCENT_ALT, line_width=1, transparency=20)
    add_line(slide, Inches(11.1), Inches(0.7), Inches(0.75), Inches(0.02), ACCENT)
    add_line(slide, Inches(11.1), Inches(0.92), Inches(0.48), Inches(0.02), ACCENT_ALT)


def add_image_placeholder(slide, left, top, width, height, label="Replace with image"):
    frame = add_rect(slide, left, top, width, height, PANEL, line_color=ACCENT_ALT, line_width=1, transparency=5)
    add_textbox(
        slide,
        left + Inches(0.2),
        top + height / 2 - Inches(0.15),
        width - Inches(0.4),
        Inches(0.3),
        label,
        18,
        color=MUTED,
        font_name=BODY_FONT,
        align=PP_ALIGN.CENTER,
    )
    return frame


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(12.15), Inches(6.15), PANEL, line_color=LINE, transparency=12)
    add_line(slide, Inches(0.9), Inches(1.05), Inches(3.0), Inches(0.05), ACCENT)
    add_textbox(slide, Inches(0.95), Inches(1.35), Inches(7.5), Inches(1.0), "Presentation Title", 30, font_name=TITLE_FONT, bold=True)
    add_textbox(slide, Inches(0.98), Inches(2.2), Inches(6.5), Inches(0.7), "Subtitle or project focus", 18, color=MUTED)
    add_textbox(slide, Inches(0.98), Inches(5.55), Inches(5.2), Inches(0.25), "Your Name", 14, color=TEXT, bold=True)
    add_textbox(slide, Inches(0.98), Inches(5.88), Inches(5.2), Inches(0.25), "Course / Teacher / Date", 11, color=MUTED)
    add_rect(slide, Inches(8.95), Inches(1.15), Inches(2.9), Inches(4.8), PANEL_ALT, line_color=ACCENT_ALT, transparency=10)
    add_line(slide, Inches(9.22), Inches(1.5), Inches(2.1), Inches(0.02), ACCENT_ALT)
    add_line(slide, Inches(9.22), Inches(1.78), Inches(1.45), Inches(0.02), ACCENT)
    add_textbox(slide, Inches(9.25), Inches(2.25), Inches(2.2), Inches(0.5), "Template", 20, font_name=TITLE_FONT, bold=True, color=TEXT)
    add_textbox(slide, Inches(9.25), Inches(2.85), Inches(2.1), Inches(2.2), "Reusable\nSchool-friendly\nDark fantasy theme", 18, color=MUTED)
    add_footer(slide)


def create_section_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_rect(slide, Inches(0.8), Inches(1.25), Inches(11.7), Inches(4.8), PANEL, line_color=LINE, transparency=10)
    add_line(slide, Inches(1.15), Inches(2.2), Inches(1.9), Inches(0.05), ACCENT)
    add_textbox(slide, Inches(1.15), Inches(2.45), Inches(8.0), Inches(0.8), "Section Title", 28, font_name=TITLE_FONT, bold=True)
    add_textbox(slide, Inches(1.18), Inches(3.15), Inches(7.0), Inches(0.55), "Short divider statement or transition", 16, color=MUTED)
    add_corner_accent(slide)
    add_footer(slide)


def create_standard_content(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Slide Title", "Use for standard explanations, definitions, or key points.")
    add_rect(slide, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.75), PANEL, line_color=LINE, transparency=8)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.0),
        Inches(10.8),
        Inches(3.8),
        "Key Points",
        [
            "First main point goes here",
            "Second point with room for explanation",
            "Third point to support the topic",
            "Optional final takeaway or example",
        ],
    )
    add_footer(slide)


def create_two_column(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Two-Column Layout", "Useful for splitting ideas, evidence, or explanation and example.")
    add_rect(slide, Inches(0.7), Inches(1.7), Inches(5.75), Inches(4.75), PANEL, line_color=LINE, transparency=8)
    add_rect(slide, Inches(6.95), Inches(1.7), Inches(5.75), Inches(4.75), PANEL, line_color=LINE, transparency=8)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.0),
        Inches(5.1),
        Inches(3.8),
        "Left Column",
        ["Main idea", "Supporting detail", "Example or note"],
    )
    add_bullets(
        slide,
        Inches(7.25),
        Inches(2.0),
        Inches(5.1),
        Inches(3.8),
        "Right Column",
        ["Comparison point", "Supporting detail", "Example or note"],
    )
    add_footer(slide)


def create_image_focused(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Image-Focused Slide", "Designed for screenshots, diagrams, maps, or figures.")
    add_image_placeholder(slide, Inches(0.8), Inches(1.75), Inches(8.6), Inches(4.9), "Replace with image / chart / diagram")
    add_rect(slide, Inches(9.7), Inches(1.75), Inches(2.65), Inches(4.9), PANEL, line_color=LINE, transparency=8)
    add_textbox(slide, Inches(9.95), Inches(2.05), Inches(2.15), Inches(0.5), "Caption", 20, font_name=BODY_FONT, bold=True)
    add_textbox(
        slide,
        Inches(9.95),
        Inches(2.55),
        Inches(2.0),
        Inches(2.7),
        "Short explanation of what the image shows and why it matters.",
        16,
        color=MUTED,
    )
    add_textbox(slide, Inches(9.95), Inches(5.45), Inches(2.0), Inches(0.35), "Source / Figure note", 11, color=MUTED)
    add_footer(slide)


def create_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Comparison Slide", "Use for pros/cons, before/after, or two competing ideas.")
    add_rect(slide, Inches(0.9), Inches(2.0), Inches(5.3), Inches(3.95), PANEL, line_color=ACCENT_ALT, transparency=6)
    add_rect(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(3.95), PANEL, line_color=ACCENT, transparency=6)
    add_textbox(slide, Inches(1.2), Inches(2.28), Inches(2.6), Inches(0.45), "Option A", 22, font_name=TITLE_FONT, bold=True)
    add_textbox(slide, Inches(7.4), Inches(2.28), Inches(2.6), Inches(0.45), "Option B", 22, font_name=TITLE_FONT, bold=True)
    add_bullets(slide, Inches(1.2), Inches(2.8), Inches(4.5), Inches(2.7), "", ["Point one", "Point two", "Point three"])
    add_bullets(slide, Inches(7.4), Inches(2.8), Inches(4.5), Inches(2.7), "", ["Point one", "Point two", "Point three"])
    add_footer(slide)


def create_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Timeline / Process", "Best for steps, workflows, or chronological events.")
    add_line(slide, Inches(1.2), Inches(3.55), Inches(10.8), Inches(0.03), LINE)

    step_x = [1.0, 4.0, 7.0, 10.0]
    labels = ["Step 01", "Step 02", "Step 03", "Step 04"]
    desc = [
        "Introduce the first stage",
        "Explain the next action",
        "Show development or result",
        "End with outcome or reflection",
    ]

    for i, x in enumerate(step_x):
        add_rect(slide, Inches(x), Inches(2.8), Inches(1.4), Inches(1.4), PANEL_ALT if i % 2 == 0 else PANEL, line_color=ACCENT if i % 2 else ACCENT_ALT, transparency=4)
        add_textbox(slide, Inches(x), Inches(3.2), Inches(1.4), Inches(0.35), f"{i + 1}", 22, font_name=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.05), Inches(4.4), Inches(1.5), Inches(0.3), labels[i], 14, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.45), Inches(4.8), Inches(2.3), Inches(0.9), desc[i], 12, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide)


def create_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Conclusion", "Wrap up the main idea and end with a clear final message.")
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(7.0), Inches(4.7), PANEL, line_color=LINE, transparency=8)
    add_textbox(slide, Inches(1.1), Inches(2.2), Inches(6.2), Inches(0.5), "Final Takeaway", 24, font_name=TITLE_FONT, bold=True)
    add_bullets(
        slide,
        Inches(1.1),
        Inches(2.85),
        Inches(6.0),
        Inches(2.8),
        "",
        ["Restate the most important idea", "Mention one supporting insight", "Close with a clear ending statement"],
    )
    add_rect(slide, Inches(8.2), Inches(1.8), Inches(4.0), Inches(4.7), PANEL_ALT, line_color=ACCENT_ALT, transparency=10)
    add_textbox(slide, Inches(8.55), Inches(2.35), Inches(3.2), Inches(0.8), "Thank You", 28, font_name=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.65), Inches(3.25), Inches(3.0), Inches(1.1), "Questions\nor discussion", 20, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_section_divider(prs)
    create_standard_content(prs)
    create_two_column(prs)
    create_image_focused(prs)
    create_comparison(prs)
    create_timeline(prs)
    create_conclusion(prs)

    prs.save(OUT_PATH)
    print(f"Created {OUT_PATH}")


if __name__ == "__main__":
    main()
