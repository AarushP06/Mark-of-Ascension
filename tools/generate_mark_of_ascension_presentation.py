from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from generate_presentation_template import (
    ACCENT,
    ACCENT_ALT,
    BODY_FONT,
    MUTED,
    PANEL,
    PANEL_ALT,
    TEXT,
    TITLE_FONT,
    add_bullets,
    add_footer,
    add_header,
    add_image_placeholder,
    add_line,
    add_rect,
    add_textbox,
    set_background,
)


OUT_PATH = "Mark_of_Ascension_Presentation.pptx"
IMAGE_PATH = Path("Assets/ChatGPT Image May 13, 2026, 05_47_32 PM.png")


def add_cover_image(slide, path, left, top, width, height, transparency_panel=True):
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if transparency_panel:
        overlay = add_rect(slide, left, top, width, height, PANEL, line_color=PANEL, transparency=40)
        overlay.line.fill.background()


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    if IMAGE_PATH.exists():
        add_cover_image(slide, IMAGE_PATH, Inches(0), Inches(0), prs.slide_width, prs.slide_height)

    add_rect(slide, Inches(0.7), Inches(0.8), Inches(6.7), Inches(5.7), PANEL, line_color=ACCENT_ALT, transparency=18)
    add_line(slide, Inches(1.0), Inches(1.25), Inches(2.8), Inches(0.05), ACCENT)
    add_textbox(slide, Inches(1.05), Inches(1.55), Inches(5.9), Inches(1.1), "Mark of Ascension", 30, font_name=TITLE_FONT, bold=True)
    add_textbox(slide, Inches(1.05), Inches(2.45), Inches(5.6), Inches(0.9), "A 2D top-down dark fantasy action game built in Unity", 19, color=MUTED)
    add_textbox(slide, Inches(1.05), Inches(4.95), Inches(3.5), Inches(0.3), "Game Development Project", 14, color=TEXT, bold=True)
    add_textbox(slide, Inches(1.05), Inches(5.3), Inches(4.2), Inches(0.35), "Aarush Patel", 12, color=MUTED)
    add_textbox(slide, Inches(1.05), Inches(5.58), Inches(4.6), Inches(0.35), "Built from the current Unity project repository", 12, color=MUTED)
    add_footer(slide, "Mark of Ascension", "Game Development Presentation")


def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Project Overview", "High-level concept and scope based on the current repo.")
    add_rect(slide, Inches(0.7), Inches(1.75), Inches(12.0), Inches(4.75), PANEL, transparency=8)

    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.05),
        Inches(10.8),
        Inches(3.9),
        "What The Game Is",
        [
            "A 2D top-down action game with a dark fantasy setting and stage-based progression.",
            "The player clears enemy waves, defeats bosses, and unlocks stronger abilities between stages.",
            "The project includes a full playable loop from main menu to final-stage victory.",
            "Built in Unity with separate scene, gameplay, UI, progression, and audio systems.",
        ],
    )
    add_footer(slide, "Overview", "Source: README.md + Assets/Scripts")


def create_visual_theme_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Visual Direction", "The project uses a dark fantasy presentation style consistent with the game theme.")

    if IMAGE_PATH.exists():
        slide.shapes.add_picture(str(IMAGE_PATH), Inches(0.8), Inches(1.85), width=Inches(8.0), height=Inches(4.85))
    else:
        add_image_placeholder(slide, Inches(0.8), Inches(1.85), Inches(8.0), Inches(4.85))

    add_rect(slide, Inches(9.1), Inches(1.85), Inches(3.45), Inches(4.85), PANEL, transparency=8)
    add_textbox(slide, Inches(9.4), Inches(2.15), Inches(2.7), Inches(0.45), "Art Direction", 22, font_name=TITLE_FONT, bold=True)
    add_bullets(
        slide,
        Inches(9.4),
        Inches(2.7),
        Inches(2.6),
        Inches(3.2),
        "",
        [
            "Dark fantasy mood",
            "Readable top-down presentation",
            "Room for bosses, hazards, and progression feedback",
            "School-friendly without losing visual identity",
        ],
    )
    add_footer(slide, "Visual Theme", "Project art asset in Assets/")


def create_game_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Gameplay Flow", "The current progression loop is scene-based and easy to explain to players.")
    add_line(slide, Inches(1.2), Inches(3.45), Inches(10.8), Inches(0.03), ACCENT_ALT)

    steps = [
        ("Main Menu", "Start game"),
        ("Lobby", "Prepare and choose stage"),
        ("Stage01", "Intro combat and boss"),
        ("Stage02", "Stronger enemies and reward"),
        ("Stage03", "Final boss and victory"),
    ]
    x_positions = [0.95, 3.2, 5.45, 7.7, 9.95]

    for i, ((title, desc), x) in enumerate(zip(steps, x_positions)):
        add_rect(slide, Inches(x), Inches(2.75), Inches(1.5), Inches(1.35), PANEL_ALT if i % 2 == 0 else PANEL, line_color=ACCENT if i % 2 == 0 else ACCENT_ALT, transparency=5)
        add_textbox(slide, Inches(x + 0.1), Inches(3.0), Inches(1.3), Inches(0.35), str(i + 1), 20, font_name=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.15), Inches(4.25), Inches(1.8), Inches(0.35), title, 13, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.3), Inches(4.62), Inches(2.1), Inches(0.7), desc, 11, color=MUTED, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        Inches(0.95),
        Inches(5.85),
        Inches(11.0),
        Inches(0.45),
        "Core loop: clear enemies -> unlock boss path -> defeat boss -> earn reward -> use portal to continue or return.",
        14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Flow", "MainMenu -> SC_Lobby -> Stage01 -> Stage02 -> Stage03")


def create_stage_rewards_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Stage Progression and Rewards", "Each stage adds difficulty and gives a meaningful gameplay upgrade.")

    headers = ["Stage01", "Stage02", "Stage03"]
    bodies = [
        ["Early dungeon introduction", "Wave clear unlocks boss path", "Reward: Poison Strike", "Bonus: Max Health +2"],
        ["Tougher fortress stage", "Stronger mid-game combat", "Reward: Flame Strike", "Bonus: damage up + faster attacks"],
        ["Dark fortress finale", "Hazards + final boss pressure", "Reward: victory state", "Outcome: run complete message"],
    ]
    xs = [0.85, 4.38, 7.92]

    for i, x in enumerate(xs):
        add_rect(slide, Inches(x), Inches(2.0), Inches(3.0), Inches(4.4), PANEL if i != 1 else PANEL_ALT, line_color=ACCENT if i != 1 else ACCENT_ALT, transparency=8)
        add_textbox(slide, Inches(x + 0.22), Inches(2.25), Inches(2.5), Inches(0.4), headers[i], 22, font_name=TITLE_FONT, bold=True)
        add_bullets(slide, Inches(x + 0.22), Inches(2.8), Inches(2.5), Inches(3.0), "", bodies[i])

    add_footer(slide, "Stages & Rewards", "Source: README.md + PlayerProgression.cs")


def create_systems_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Core Gameplay Systems", "The project combines player, combat, progression, UI, and feedback systems.")
    add_rect(slide, Inches(0.7), Inches(1.8), Inches(5.85), Inches(4.65), PANEL, transparency=8)
    add_rect(slide, Inches(6.85), Inches(1.8), Inches(5.85), Inches(4.65), PANEL, transparency=8)

    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.1),
        Inches(5.1),
        Inches(3.9),
        "Gameplay Systems",
        [
            "Persistent player spawning across scenes",
            "Top-down movement and melee attack",
            "Enemy AI with contact damage",
            "Boss spawning after enemy-wave clear",
            "Hazard zones that punish positioning",
        ],
    )
    add_bullets(
        slide,
        Inches(7.15),
        Inches(2.1),
        Inches(5.1),
        Inches(3.9),
        "Feedback Systems",
        [
            "Level-up reward notifications between stages",
            "Poison and fire damage-over-time effects",
            "Audio for combat, boss events, portals, and lobby music",
            "Game-over flow with return to main menu",
            "Final victory message after Stage03",
        ],
    )
    add_footer(slide, "Systems", "Gameplay + UI + Audio")


def create_technical_structure_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Technical Structure", "The repo is organized around scenes, gameplay scripts, and UI systems.")
    add_rect(slide, Inches(0.75), Inches(1.85), Inches(5.75), Inches(4.7), PANEL, transparency=8)
    add_rect(slide, Inches(6.85), Inches(1.85), Inches(5.75), Inches(4.7), PANEL, transparency=8)

    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.1),
        Inches(5.1),
        Inches(3.8),
        "Important Folders",
        [
            "Assets/Scenes: MainMenu, SC_Lobby, Stage01, Stage02, Stage03",
            "Assets/Scripts/Gameplay: combat, AI, portals, hazards, progression",
            "Assets/Scripts/UI: menu flow, health bar, reward notifications",
            "Assets/Resources/Audio: reusable sound effects and music clips",
        ],
    )
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.1),
        Inches(5.1),
        Inches(3.8),
        "Key Scripts",
        [
            "PlayerProgression.cs manages rewards and victory states",
            "Stage01FlowController.cs handles enemy waves, boss spawn, and portal unlock",
            "ScenePortal.cs manages scene transitions",
            "GameAudio.cs centralizes audio playback",
        ],
    )
    add_footer(slide, "Structure", "Source: Assets/Scenes + Assets/Scripts")


def create_progression_logic_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "How The Progression Loop Works", "The gameplay loop is implemented as a clear chain of events inside the project scripts.")
    add_line(slide, Inches(1.15), Inches(3.5), Inches(10.8), Inches(0.03), ACCENT_ALT)

    steps = [
        ("Enemy wave", "Small enemies are tracked"),
        ("Boss unlock", "Path blockers are removed"),
        ("Boss spawn", "Boss is instantiated and configured"),
        ("Reward", "PlayerProgression grants upgrade"),
        ("Portal exit", "Scene portal or victory message"),
    ]
    x_positions = [0.95, 3.15, 5.35, 7.55, 9.75]

    for i, ((title, desc), x) in enumerate(zip(steps, x_positions)):
        add_rect(slide, Inches(x), Inches(2.75), Inches(1.45), Inches(1.45), PANEL_ALT if i in (0, 2, 4) else PANEL, line_color=ACCENT if i in (0, 2, 4) else ACCENT_ALT, transparency=4)
        add_textbox(slide, Inches(x + 0.08), Inches(3.06), Inches(1.28), Inches(0.3), title, 13, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.22), Inches(4.45), Inches(1.9), Inches(0.85), desc, 11, color=MUTED, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        Inches(1.0),
        Inches(5.8),
        Inches(11.0),
        Inches(0.55),
        "This structure makes the project easy to extend with additional stages, rewards, and boss encounters later.",
        14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Implementation", "Stage01FlowController.cs + PlayerProgression.cs")


def create_current_status_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Current Status and Strengths", "What is already working in the project right now.")
    add_rect(slide, Inches(0.8), Inches(1.85), Inches(11.9), Inches(4.75), PANEL, transparency=8)
    add_bullets(
        slide,
        Inches(1.1),
        Inches(2.15),
        Inches(10.9),
        Inches(3.9),
        "Playable Features",
        [
            "Full flow from main menu to final stage is implemented.",
            "Boss rewards directly change combat by unlocking poison and fire effects.",
            "UI gives feedback through health bars, reward popups, death flow, and final victory messaging.",
            "Audio support now covers combat, boss events, portal actions, and lobby music.",
            "The codebase is organized enough to keep expanding with more stages or polish.",
        ],
    )
    add_footer(slide, "Status", "Current repo build")


def create_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Conclusion", "Mark of Ascension already demonstrates a complete gameplay loop and clear progression system.")
    add_rect(slide, Inches(0.8), Inches(1.9), Inches(7.05), Inches(4.55), PANEL, transparency=8)
    add_rect(slide, Inches(8.15), Inches(1.9), Inches(4.05), Inches(4.55), PANEL_ALT, transparency=10)

    add_bullets(
        slide,
        Inches(1.1),
        Inches(2.2),
        Inches(6.25),
        Inches(3.6),
        "Key Takeaways",
        [
            "The project combines gameplay, progression, UI, and audio into one coherent loop.",
            "Its stage-based structure makes the design easy to explain and easy to expand.",
            "The current repo already supports a strong class presentation because the systems are visible in both scenes and code.",
        ],
    )
    add_textbox(slide, Inches(8.5), Inches(2.45), Inches(3.2), Inches(0.6), "Thank You", 28, font_name=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.55), Inches(3.35), Inches(3.0), Inches(1.0), "Questions\nor discussion", 20, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Conclusion", "Mark of Ascension")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_overview_slide(prs)
    create_visual_theme_slide(prs)
    create_game_flow_slide(prs)
    create_stage_rewards_slide(prs)
    create_systems_slide(prs)
    create_technical_structure_slide(prs)
    create_progression_logic_slide(prs)
    create_current_status_slide(prs)
    create_conclusion_slide(prs)

    prs.save(OUT_PATH)
    print(f"Created {OUT_PATH}")


if __name__ == "__main__":
    main()
